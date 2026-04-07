"""Office document converter: DOCX, PPTX, XLSX.

Primary: MarkItDown (microsoft/markitdown, MIT)
Fallback: pandoc subprocess

v0.3: adds embedded image extraction for DOCX and PPTX via python-docx /
python-pptx (optional extras). Extracted images are saved to
output_dir/images/ and included in ConverterResult.images so they feed
into the Visual Enhancement Pipeline.
"""

from __future__ import annotations

import logging
import subprocess
from pathlib import Path

from drop2md.converters import BaseConverter, ConversionError, ConverterResult

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Embedded image extraction (VEP-4)
# ---------------------------------------------------------------------------

def _extract_docx_images(path: Path, output_dir: Path) -> list[Path]:
    """Extract embedded images from a DOCX file using python-docx.

    Returns list of saved image paths; returns [] if python-docx is not
    installed or the file contains no images.
    """
    try:
        from docx import Document  # python-docx
    except ImportError:
        log.debug("python-docx not available — skipping DOCX image extraction")
        return []

    img_dir = output_dir / "images"
    img_dir.mkdir(parents=True, exist_ok=True)
    saved: list[Path] = []

    try:
        doc = Document(str(path))
        for idx, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.reltype:
                img_part = rel.target_part
                ext = Path(img_part.partname).suffix or ".png"
                dest = img_dir / f"{path.stem}_img{idx}{ext}"
                dest.write_bytes(img_part.blob)
                saved.append(dest)
                log.debug("Extracted DOCX image: %s", dest.name)
    except Exception as exc:
        log.debug("DOCX image extraction error for %s: %s", path.name, exc)

    return saved


def _extract_pptx_images(path: Path, output_dir: Path) -> list[Path]:
    """Extract embedded images from a PPTX file using python-pptx.

    Returns list of saved image paths; returns [] if python-pptx is not
    installed or the file contains no images.
    """
    try:
        from pptx import Presentation  # python-pptx
        from pptx.util import Inches  # noqa: F401  (confirms package is functional)
    except ImportError:
        log.debug("python-pptx not available — skipping PPTX image extraction")
        return []

    img_dir = output_dir / "images"
    img_dir.mkdir(parents=True, exist_ok=True)
    saved: list[Path] = []

    try:
        prs = Presentation(str(path))
        idx = 0
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        img_bytes = shape.image.blob
                        ext = f".{shape.image.ext}"
                        dest = img_dir / f"{path.stem}_slide{slide_num}_img{idx}{ext}"
                        dest.write_bytes(img_bytes)
                        saved.append(dest)
                        log.debug("Extracted PPTX image: %s", dest.name)
                        idx += 1
                    except Exception as exc:
                        log.debug("Could not extract PPTX shape image: %s", exc)
    except Exception as exc:
        log.debug("PPTX image extraction error for %s: %s", path.name, exc)

    return saved


def _extract_office_images(path: Path, output_dir: Path) -> list[Path]:
    """Dispatch to the correct image extractor based on file suffix."""
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return _extract_docx_images(path, output_dir)
    if suffix == ".pptx":
        return _extract_pptx_images(path, output_dir)
    return []


# ---------------------------------------------------------------------------
# Converters
# ---------------------------------------------------------------------------

class MarkItDownConverter(BaseConverter):
    """MarkItDown converter for Office formats (DOCX, PPTX, XLSX)."""

    name = "markitdown"

    @classmethod
    def is_available(cls) -> bool:
        try:
            import markitdown  # noqa: F401

            return True
        except ImportError:
            return False

    def convert(self, path: Path, output_dir: Path) -> ConverterResult:
        from markitdown import MarkItDown

        # Extract embedded images before running MarkItDown so they are
        # available to the VEP pipeline regardless of what MarkItDown does.
        images = _extract_office_images(path, output_dir)

        md = MarkItDown()
        result = md.convert(str(path))
        return ConverterResult(
            markdown=result.text_content,
            images=images,
            converter_used=self.name,
        )


class PandocOfficeConverter(BaseConverter):
    """Pandoc subprocess fallback for DOCX files."""

    name = "pandoc"

    @classmethod
    def is_available(cls) -> bool:
        return subprocess.run(
            ["pandoc", "--version"],
            capture_output=True,
            check=False,
        ).returncode == 0

    def convert(self, path: Path, output_dir: Path) -> ConverterResult:
        suffix = path.suffix.lower()
        if suffix not in {".docx", ".odt", ".rtf"}:
            raise ConversionError(f"Pandoc fallback does not support {suffix}")

        images = _extract_office_images(path, output_dir)

        result = subprocess.run(
            ["pandoc", "-f", "docx", "-t", "gfm", str(path)],
            capture_output=True,
            text=True,
            check=False,
        )
        if result.returncode != 0:
            raise ConversionError(f"pandoc failed: {result.stderr}")
        return ConverterResult(
            markdown=result.stdout,
            images=images,
            converter_used=self.name,
        )


class OfficeConverter(BaseConverter):
    """Routes DOCX/PPTX/XLSX through MarkItDown with Pandoc fallback."""

    name = "office"

    def convert(self, path: Path, output_dir: Path) -> ConverterResult:
        _converters: list[type[BaseConverter]] = [MarkItDownConverter, PandocOfficeConverter]
        for ConverterClass in _converters:
            if not ConverterClass.is_available():
                continue
            try:
                result = ConverterClass().convert(path, output_dir)
                log.info("Converted %s via %s", path.name, ConverterClass.name)
                return result
            except Exception as exc:
                log.warning("%s failed for %s: %s", ConverterClass.name, path.name, exc)

        raise ConversionError(f"No office converter available for {path}")
